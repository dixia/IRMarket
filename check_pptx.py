from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

prs = Presentation('docs/Monoracle.pptx')
print(f'Slide count: {len(prs.slides)}')
print(f'Slide width: {prs.slide_width/914400}in, height: {prs.slide_height/914400}in')

for i, slide in enumerate(prs.slides):
    print(f'\n=== Slide {i+1} ===')
    print(f'Layout: {slide.slide_layout.name}')
    for shape in slide.shapes:
        print(f'  Shape: {shape.shape_type}, name={shape.name}, pos=({shape.left},{shape.top}), size=({shape.width},{shape.height})')
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    font = run.font
                    try:
                        color = font.color.rgb if font.color and font.color.type else None
                    except:
                        color = 'theme/inherited'
                    print(f'    Text: "{run.text[:100]}" size={font.size} bold={font.bold} color={color}')
        # Check for fill color
        if hasattr(shape, 'fill'):
            try:
                if shape.fill.type is not None:
                    print(f'    Fill type: {shape.fill.type}')
                    if shape.fill.type == 1:  # solid
                        print(f'    Fill color: {shape.fill.fore_color.rgb}')
            except:
                pass
